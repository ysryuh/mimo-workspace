#!/usr/bin/env python3
"""AMMR-H600 통합설계서 PPT 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Colors ───
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x00, 0x96, 0xD6)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE0, 0xE0, 0xE0)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
SECTION_BG = RGBColor(0x0A, 0x0A, 0x1A)
CARD_BG = RGBColor(0x25, 0x25, 0x40)
GREEN = RGBColor(0x00, 0xC8, 0x53)
RED = RGBColor(0xFF, 0x45, 0x45)
YELLOW = RGBColor(0xFF, 0xD6, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color=None, line_color=None, shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, left, top, width, height)
    shp.fill.background()
    if fill_color:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="맑은 고딕"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_text(slide, left, top, width, height, items, font_size=16, color=WHITE, bullet_color=ACCENT_BLUE, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"▸ {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "맑은 고딕"
        p.space_after = spacing
    return txBox


def add_card(slide, left, top, width, height, title, items, title_color=ACCENT_BLUE):
    """Add a card with title and bullet items"""
    card = add_shape(slide, left, top, width, height, fill_color=CARD_BG)
    card.shadow.inherit = False
    # Title
    add_text(slide, left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.4),
             title, font_size=16, color=title_color, bold=True)
    # Items
    y = top + Inches(0.5)
    for item in items:
        add_text(slide, left + Inches(0.25), y, width - Inches(0.5), Inches(0.3),
                 f"• {item}", font_size=13, color=LIGHT_GRAY)
        y += Inches(0.28)


def add_divider(slide, left, top, width, color=ACCENT_BLUE):
    shp = add_shape(slide, left, top, width, Pt(3), fill_color=color)


# ═══════════════════════════════════════════════════════
# SLIDE 1: 표지
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, SECTION_BG)

# Accent line
add_shape(slide, Inches(1.5), Inches(2.2), Inches(2), Pt(4), fill_color=ACCENT_BLUE)

add_text(slide, Inches(1.5), Inches(2.4), Inches(10), Inches(1.2),
         "고하중 AMMR 통합 설계서", font_size=44, color=WHITE, bold=True)
add_text(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.6),
         "High-Load Autonomous Mobile Manipulator Robot — Integrated Design Document", font_size=20, color=LIGHT_GRAY)
add_text(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.4),
         "AMMR-H600  |  v0.1 Draft  |  2026-04-01", font_size=16, color=MID_GRAY)
add_text(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.4),
         "제임스 (James) + 미모 (Mimo)", font_size=16, color=MID_GRAY)

# Decorative shape
add_shape(slide, Inches(10), Inches(1), Inches(2.5), Inches(5.5), fill_color=CARD_BG)
add_text(slide, Inches(10.3), Inches(2), Inches(2), Inches(3),
         "🤖\n\n양팔\n30kg+\n600kg AMR\n자율주행", font_size=18, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# SLIDE 2: 목차
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.6),
         "목차", font_size=32, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1), Inches(1.5))

toc_items = [
    ("01", "제품 개요"),
    ("02", "시장 분석 및 기회"),
    ("03", "타겟 요구사항"),
    ("04", "경쟁사 분석"),
    ("05", "시스템 아키텍처"),
    ("06", "하위 시스템 설계"),
    ("07", "소프트웨어 아키텍처"),
    ("08", "안전 및 규격"),
    ("09", "개발 로드맵"),
    ("10", "리스크 분석"),
]

for i, (num, title) in enumerate(toc_items):
    col = i % 2
    row = i // 2
    x = Inches(1) + col * Inches(5.5)
    y = Inches(1.5) + row * Inches(1.05)
    add_text(slide, x, y, Inches(0.6), Inches(0.5), num, font_size=24, color=ACCENT_BLUE, bold=True)
    add_text(slide, x + Inches(0.7), y + Inches(0.05), Inches(4), Inches(0.4), title, font_size=18, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════
# SLIDE 3: 제품 개요
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
         "01  제품 개요", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
         "AMMR-H600: 세계 최초 고하중 통합 AMMR", font_size=30, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.5), Inches(2))

# Key value proposition
add_shape(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(1.5), fill_color=CARD_BG)
add_text(slide, Inches(1.1), Inches(1.95), Inches(11), Inches(1.2),
         '"양팔 + 고가반 + 고적재 AMR"을 하나의 통합 플랫폼으로\n'
         '→ 공간 효율 ↑, 재배치 자유도 ↑, 총비용 ↓',
         font_size=20, color=LIGHT_GRAY)

# 4 key specs
specs = [
    ("🦾 양팔", "좌우 독립/협동 작업", ACCENT_BLUE),
    ("⚖️ 가반 30kg+", "각 팔 30kg, 합산 60kg+", ACCENT_ORANGE),
    ("🚛 적재 600kg+", "고하중 AMR 플랫폼", GREEN),
    ("🧭 자율주행", "SLAM 기반, 인프라 무수정", YELLOW),
]

for i, (icon_title, desc, color) in enumerate(specs):
    x = Inches(0.8) + i * Inches(3)
    add_shape(slide, x, Inches(3.8), Inches(2.7), Inches(2.8), fill_color=CARD_BG)
    add_text(slide, x + Inches(0.2), Inches(4.0), Inches(2.3), Inches(0.5),
             icon_title, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), Inches(4.6), Inches(2.3), Inches(1.5),
             desc, font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Blue ocean callout
add_shape(slide, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.4), fill_color=ACCENT_BLUE)
add_text(slide, Inches(1), Inches(6.9), Inches(11), Inches(0.4),
         '★ "양팔 + 30kg+ + 600kg AMR 통합" — 시장에 존재하지 않는 제품 = 블루오션',
         font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
# SLIDE 4: 시장 분석
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
         "02  시장 분석 및 기회", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
         "AMR 시장 성장 & 블루오션 근거", font_size=30, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.5), Inches(2))

# Market size card
add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5),
         "📈 시장 트렌드", [
             "글로벌 AMR 시장: 2025년 ~$5.5B",
             "2030년 $12B+ (CAGR ~17%)",
             "모바일 매니퓰레이터: 초기 단계",
             "스마트 팩토리 수요 급증",
         ])

# Target personas
add_card(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(2.5),
         "🎯 타겟 고객", [
             "1차: 대형 제조 공장 (반도체/디스플레이/자동차)",
             "기존 AGV + 고정 로봇 운영 중",
             "바닥면적 제약, 라인 변경 빈번",
             "2차: 물류 허브 (대형 팔레트 처리)",
         ])

# Blue ocean matrix
add_shape(slide, Inches(0.8), Inches(4.6), Inches(11.5), Inches(2.6), fill_color=CARD_BG)
add_text(slide, Inches(1.1), Inches(4.7), Inches(11), Inches(0.4),
         "🔍 블루오션 조합 분석", font_size=18, color=ACCENT_ORANGE, bold=True)

matrix_items = [
    ("양팔 매니퓰레이터", "✅ 있음 (ABB YuMi, Kawasaki duAro)"),
    ("30kg+ 가반 양팔", "⚠️ 제한적 (단팔 30kg+는 다수)"),
    ("600kg+ AMR", "✅ 있음 (MiR600, OTTO 1500)"),
    ("양팔 + 30kg+ + 600kg AMR 통합", "❌ 시장에 없음 ← 우리 타겟"),
]
for i, (label, status) in enumerate(matrix_items):
    y = Inches(5.2) + i * Inches(0.48)
    add_text(slide, Inches(1.3), y, Inches(4), Inches(0.4), label, font_size=14, color=LIGHT_GRAY)
    color = GREEN if "✅" in status else (YELLOW if "⚠️" in status else RED)
    add_text(slide, Inches(6), y, Inches(6), Inches(0.4), status, font_size=14, color=color)


# ═══════════════════════════════════════════════════════
# SLIDE 5: 타겟 요구사항
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
         "03  타겟 요구사항", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
         "핵심 성능 사양", font_size=30, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.5), Inches(2))

specs_data = [
    ("AMR 적재능력", "≥ 600kg", ACCENT_BLUE),
    ("팔 가반중량 (각)", "≥ 30kg", ACCENT_BLUE),
    ("팔 작업반경", "≥ 1,500mm", ACCENT_ORANGE),
    ("자율주행", "SLAM 기반", GREEN),
    ("최대 이동속도", "1.5 m/s (하중)", ACCENT_ORANGE),
    ("위치 정밀도 (AMR)", "±10mm", MID_GRAY),
    ("위치 정밀도 (팔)", "±0.1mm", MID_GRAY),
    ("배터리 수명", "≥ 6시간", GREEN),
    ("급속 충전", "80% in 30min", GREEN),
    ("등판각", "≥ 5° (600kg)", ACCENT_ORANGE),
    ("급정지 거리", "≤ 1m", RED),
    ("통신", "Wi-Fi 6E + 5G", ACCENT_BLUE),
]

for i, (label, value, color) in enumerate(specs_data):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * Inches(4.1)
    y = Inches(1.8) + row * Inches(1.3)
    add_shape(slide, x, y, Inches(3.8), Inches(1.1), fill_color=CARD_BG)
    add_text(slide, x + Inches(0.2), y + Inches(0.1), Inches(3.4), Inches(0.3),
             label, font_size=13, color=MID_GRAY)
    add_text(slide, x + Inches(0.2), y + Inches(0.45), Inches(3.4), Inches(0.5),
             value, font_size=22, color=color, bold=True)


# ═══════════════════════════════════════════════════════
# SLIDE 6: 경쟁사 분석
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
         "04  경쟁사 분석", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
         "경쟁 매트릭스", font_size=30, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.5), Inches(2))

# Table header
headers = ["업체", "제품", "AMR 적재", "팔 가반", "양팔", "통합"]
header_widths = [Inches(2), Inches(2.2), Inches(1.5), Inches(1.5), Inches(1.2), Inches(1.5)]

x_start = Inches(0.8)
y_header = Inches(1.8)
# Header bg
add_shape(slide, x_start, y_header, Inches(10), Inches(0.45), fill_color=ACCENT_BLUE)
x = x_start
for j, (h, w) in enumerate(zip(headers, header_widths)):
    add_text(slide, x + Inches(0.1), y_header + Inches(0.05), w, Inches(0.35),
             h, font_size=13, color=WHITE, bold=True)
    x += w

# Table rows
competitors = [
    ("MiR + UR", "MC600", "600kg", "30kg (UR30)", "❌ 단팔", "반통합"),
    ("ROKAE", "CMR + CR35", "500kg", "45kg", "❌ 별도", "미통합"),
    ("KUKA", "KMP 600P", "600kg", "-", "-", "AMR만"),
    ("OTTO Motors", "OTTO 1500", "1500kg", "-", "-", "AMR만"),
    ("FANUC", "M-2000iA", "-", "900kg", "❌ 단팔", "팔만"),
    ("Rainbow", "RB-Y1", "경량", "3kg", "✅ 양팔", "통합"),
]

for i, row_data in enumerate(competitors):
    y = Inches(2.3) + i * Inches(0.5)
    bg = CARD_BG if i % 2 == 0 else DARK_BG
    add_shape(slide, x_start, y, Inches(10), Inches(0.48), fill_color=bg)
    x = x_start
    for j, (cell, w) in enumerate(zip(row_data, header_widths)):
        color = GREEN if "✅" in cell else (RED if "❌" in cell else LIGHT_GRAY)
        add_text(slide, x + Inches(0.1), y + Inches(0.08), w, Inches(0.32),
                 cell, font_size=12, color=color if ("✅" in cell or "❌" in cell) else LIGHT_GRAY)
        x += w

# Our row (highlighted)
y_ours = Inches(2.3) + 6 * Inches(0.5)
add_shape(slide, x_start, y_ours, Inches(10), Inches(0.55), fill_color=ACCENT_BLUE)
x = x_start
our_data = ["AMMR-H600", "우리", "600kg", "30kg+", "✅ 양팔", "✅ 완전통합"]
for cell, w in zip(our_data, header_widths):
    add_text(slide, x + Inches(0.1), y_ours + Inches(0.1), w, Inches(0.35),
             cell, font_size=13, color=WHITE, bold=True)
    x += w

# ROKAE callout
add_shape(slide, Inches(0.8), Inches(6.1), Inches(10), Inches(1.1), fill_color=CARD_BG)
add_text(slide, Inches(1.1), Inches(6.2), Inches(9.5), Inches(0.3),
         "⚠️ 최근접 경쟁자: ROKAE (珞石)", font_size=16, color=ACCENT_ORANGE, bold=True)
add_text(slide, Inches(1.1), Inches(6.55), Inches(9.5), Inches(0.5),
         "CR35 (45kg 가반) + CMR (500kg AMR) — 구성요소 보유, 통합 제품 미출시\n"
         "→ 600kg+ 자체 개발하면 확실한 차별화",
         font_size=13, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════
# SLIDE 7: 시스템 아키텍처
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
         "05  시스템 아키텍처", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
         "전체 시스템 구성", font_size=30, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.5), Inches(2))

# Three main subsystems
subsystems = [
    ("양팔 매니퓰레이터", [
        "Left Arm (30kg+ / 7DOF)",
        "Right Arm (30kg+ / 7DOF)",
        "엔드 이펙터 (퀵체인지)",
        "F/T 센서 (각 팔)",
        "비전 카메라 (팔 끝단)",
        "자기충돌검지 시스템",
    ], ACCENT_BLUE),
    ("AMR 플랫폼", [
        "4륜 독립구동 (Mecanum)",
        "BLDC 서보모터 × 4",
        "LiFePO4 배터리 2.5kWh+",
        "적재능력 600kg+",
        "급속 충전 (도킹 자동)",
        "섀시: ~1200×800mm",
    ], ACCENT_ORANGE),
    ("통합 제어 시스템", [
        "Motion Planner",
        "Task Scheduler",
        "Safety Manager (SIL2)",
        "Navigation Stack",
        "Sensor Fusion",
        "Fleet Manager 연동",
    ], GREEN),
]

for i, (title, items, color) in enumerate(subsystems):
    x = Inches(0.8) + i * Inches(4.1)
    add_shape(slide, x, Inches(1.8), Inches(3.8), Inches(4.5), fill_color=CARD_BG)
    # Top accent bar
    add_shape(slide, x, Inches(1.8), Inches(3.8), Pt(4), fill_color=color)
    add_text(slide, x + Inches(0.2), Inches(2.0), Inches(3.4), Inches(0.4),
             title, font_size=18, color=color, bold=True)
    for j, item in enumerate(items):
        add_text(slide, x + Inches(0.3), Inches(2.5) + j * Inches(0.36), Inches(3.2), Inches(0.32),
                 f"▸ {item}", font_size=13, color=LIGHT_GRAY)

# Communication bar
add_shape(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.7), fill_color=CARD_BG)
add_text(slide, Inches(1.1), Inches(6.6), Inches(11), Inches(0.4),
         "🔗 통신:  Wi-Fi 6E  |  5G (옵션)  |  EtherCAT (내부)  |  CAN (센서)  |  MQTT/gRPC (Fleet)",
         font_size=14, color=MID_GRAY)


# ═══════════════════════════════════════════════════════
# SLIDE 8: AMR 플랫폼 상세
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
         "06  하위 시스템 — AMR 플랫폼", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
         "구동 · 섀시 · 배터리", font_size=30, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.5), Inches(2))

# Drive system
add_card(slide, Inches(0.8), Inches(1.8), Inches(3.6), Inches(3),
         "🔧 구동 시스템", [
             "Mecanum 4륜 독립구동",
             "BLDC 서보모터 × 500W+",
             "하모닉 드라이브 감속기",
             "총 추력 3kN+",
             "최대 2.0m/s (공차)",
             "최대 1.5m/s (600kg)",
         ])

# Chassis
add_card(slide, Inches(4.8), Inches(1.8), Inches(3.6), Inches(3),
         "🏗️ 섀시 설계", [
             "강성 알루미늄 프레임",
             "팔 마운트 통합 구조",
             "치수: ~1200×800×400mm",
             "자중 목표: ≤ 250kg",
             "등판각: ≥ 5° (600kg)",
             "급정지: ≤ 1m (최대속도)",
         ])

# Battery
add_card(slide, Inches(8.8), Inches(1.8), Inches(3.6), Inches(3),
         "🔋 배터리 시스템", [
             "LiFePO4 (안전성 우선)",
             "용량 ≥ 2.5kWh",
             "48V / 72V",
             "연속 운용 ≥ 6시간",
             "급속 충전: 80% in 30min",
             "자동 도킹 충전",
         ])

# Top view diagram (text-based)
add_shape(slide, Inches(0.8), Inches(5.1), Inches(11.5), Inches(2.1), fill_color=CARD_BG)
add_text(slide, Inches(1.1), Inches(5.2), Inches(11), Inches(0.4),
         "📐 상면도 (Top View)", font_size=16, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1.1), Inches(5.7), Inches(11), Inches(1.3),
         "┌─────────────────────────────────────────────┐\n"
         "│          ┌───────────────────┐               │\n"
         "│  ○ Front │  팔 마운트 영역    │  Rear ○       │\n"
         "│  Wheel   │  (강성 프레임)     │  Wheel        │\n"
         "│          │  ┌───┐   ┌───┐   │               │\n"
         "│          │  │ L │   │ R │   │               │\n"
         "│          │  │Arm│   │Arm│   │               │\n"
         "│          └──┴───┴───┴───┴───┘               │\n"
         "│  ○ Front │  배터리 + 제어기   │  Rear ○       │\n"
         "│  Wheel   │                   │  Wheel        │\n"
         "└─────────────────────────────────────────────┘",
         font_size=11, color=LIGHT_GRAY, font_name="Consolas")


# ═══════════════════════════════════════════════════════
# SLIDE 9: 양팔 매니퓰레이터 상세
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
         "06  하위 시스템 — 양팔 매니퓰레이터", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
         "팔 사양 · 배치 · 엔드이펙터", font_size=30, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.5), Inches(2))

# Left arm specs
add_card(slide, Inches(0.8), Inches(1.8), Inches(3.6), Inches(3.5),
         "🦾 각 팔 사양", [
             "DOF: 6 + 1 (그리퍼) = 7축",
             "가반중량: ≥ 30kg (최대)",
             "연속 작업: 20kg (권장)",
             "작업반경: ≥ 1,500mm",
             "반복 정밀도: ±0.1mm",
             "최대 솔도: 180°/s",
             "구동: BLDC + 감속기",
             "무게: 각 팔 ≤ 50kg",
             "프로토콜: EtherCAT",
         ])

# Gripper
add_card(slide, Inches(4.8), Inches(1.8), Inches(3.6), Inches(3.5),
         "🤏 엔드 이펙터", [
             "기본: 2핀 전기 그리퍼",
             "30kg 인양 가능",
             "옵션: 진공 그리퍼",
             "옵션: 3핀 서보 그리퍼",
             "퀵 체인지: ISO 9409-1",
             "공구 자동 교체 지원",
             "Force/Torque 센서 내장",
             "비전 카메라 (팔 끝단)",
         ])

# Dual arm coordination
add_card(slide, Inches(8.8), Inches(1.8), Inches(3.6), Inches(3.5),
         "🤝 양팔 협동", [
             "양팔 간격: ~600mm (기본)",
             "양팔 교차 작업 가능",
             "독립 동작 모드",
             "협동 동작 모드",
             "실시간 자기충돌검지",
             "충돌 방지 알고리즘",
             "Admittance 힘 제어",
             "작업 분담 자동 최적화",
         ])

# Front view diagram
add_shape(slide, Inches(0.8), Inches(5.6), Inches(11.5), Inches(1.6), fill_color=CARD_BG)
add_text(slide, Inches(1.1), Inches(5.7), Inches(11), Inches(0.3),
         "📐 정면도 (Front View)", font_size=16, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1.1), Inches(6.1), Inches(11), Inches(0.9),
         "              ┌─────┐\n"
         "         ┌────┤     ├────┐\n"
         "         │ L  │ 상판 │  R │\n"
         "         │Arm │     │Arm │\n"
         "        ╱│    │     │    │╲\n"
         "       ╱ └────┘     └────┘ ╲",
         font_size=13, color=LIGHT_GRAY, font_name="Consolas")


# ═══════════════════════════════════════════════════════
# SLIDE 10: 센서 시스템
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
         "06  하위 시스템 — 센서", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
         "센서 융합 시스템", font_size=30, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.5), Inches(2))

sensors = [
    ("3D LiDAR ×1", "SLAM, 3D 매핑, 장애물 검지", ACCENT_BLUE),
    ("2D Safety LiDAR ×2", "전/후방 안전 영역 스캐닝 (SIL2)", RED),
    ("Depth Camera ×2", "전/후방 3D 매핑, 피킹 유도", ACCENT_ORANGE),
    ("IMU ×1", "관성 항법, 자세 추정", MID_GRAY),
    ("비전 카메라 ×2", "팔 끝단, 비전 유도 작업", GREEN),
    ("F/T 센서 ×2", "각 팔, 힘 제어 & 충돌 검지", YELLOW),
    ("UWB 안테나 ×4+", "실내 정밀 위치 (±10mm)", ACCENT_BLUE),
    ("엔코더 (각 관절)", "관절 위치 피드백", MID_GRAY),
]

for i, (name, desc, color) in enumerate(sensors):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6)
    y = Inches(1.8) + row * Inches(1.15)
    add_shape(slide, x, y, Inches(5.7), Inches(1), fill_color=CARD_BG)
    add_shape(slide, x, y, Pt(5), Inches(1), fill_color=color)
    add_text(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.2), Inches(0.3),
             name, font_size=15, color=color, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.5), Inches(5.2), Inches(0.4),
             desc, font_size=12, color=LIGHT_GRAY)

# Sensor fusion
add_shape(slide, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.7), fill_color=CARD_BG)
add_text(slide, Inches(1.1), Inches(6.6), Inches(11), Inches(0.4),
         "🔄 센서 융합:  LiDAR + Depth Camera + IMU → 확장 칼만 필터(EKF) → 통합 환경 인지도",
         font_size=14, color=ACCENT_BLUE, bold=True)


# ═══════════════════════════════════════════════════════
# SLIDE 11: 소프트웨어 아키텍처
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
         "07  소프트웨어 아키텍처", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
         "ROS 2 기반 소프트웨어 스택", font_size=30, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.5), Inches(2))

# SW layers
layers = [
    ("응용 계층", "Task API  |  Mission Editor  |  Fleet Dashboard", ACCENT_ORANGE),
    ("미들웨어", "ROS 2 (Humble/Iron)  |  gRPC  |  MQTT", ACCENT_BLUE),
    ("핵심 모듈", "Navigation  |  Manipulation  |  Perception  |  Safety", GREEN),
    ("드라이버", "EtherCAT Master  |  CAN Driver  |  Sensor Drivers", MID_GRAY),
    ("OS", "Ubuntu 22.04 + PREEMPT_RT  |  Safety PLC Firmware", RED),
]

for i, (name, desc, color) in enumerate(layers):
    y = Inches(1.8) + i * Inches(0.9)
    add_shape(slide, Inches(0.8), y, Inches(2.2), Inches(0.75), fill_color=CARD_BG)
    add_shape(slide, Inches(0.8), y, Pt(5), Inches(0.75), fill_color=color)
    add_text(slide, Inches(1), y + Inches(0.05), Inches(1.8), Inches(0.3),
             name, font_size=14, color=color, bold=True)
    add_text(slide, Inches(1), y + Inches(0.4), Inches(1.8), Inches(0.3),
             "", font_size=11, color=LIGHT_GRAY)

    add_shape(slide, Inches(3.2), y, Inches(9.3), Inches(0.75), fill_color=CARD_BG)
    add_text(slide, Inches(3.4), y + Inches(0.2), Inches(9), Inches(0.35),
             desc, font_size=13, color=LIGHT_GRAY)

# Key modules
add_text(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.3),
         "주요 기술:  Cartographer (SLAM)  |  Nav2 (경로계획)  |  MoveIt 2 (Manipulation)  |  Gazebo/Isaac Sim (시뮬레이션)",
         font_size=13, color=MID_GRAY)


# ═══════════════════════════════════════════════════════
# SLIDE 12: 안전 및 규격
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
         "08  안전 및 규격", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
         "안전 시스템 아키텍처", font_size=30, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.5), Inches(2))

# Safety standards
add_card(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.3),
         "📋 적용 규격", [
             "ISO 3691-4:2023 — 무인 운반차 안전 (필수)",
             "ISO 10218-1/2 — 산업 로봇 안전 (필수)",
             "ISO/TS 15066 — 협동 로봇 안전 (필수)",
             "ISO 13849-1 — 안전 제어 PLd 이상",
             "ANSI R15.08 — 산업용 모바일 로봇",
         ])

# Safety architecture
add_card(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(2.3),
         "🛡️ 안전 아키텍처 (2채널)", [
             "Main Control (표준 채널) + Safety PLC (안전 채널)",
             "SIL2 / PLd 등급 목표",
             "E-Stop: 하드웨어 + 소프트웨어 이중",
             "Safe Torque Off (STO) 기능",
             "독립 안전 검증",
         ])

# Safety zones
add_shape(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(2.8), fill_color=CARD_BG)
add_text(slide, Inches(1.1), Inches(4.5), Inches(11), Inches(0.3),
         "🚧 안전 존 (Safety Zones)", font_size=18, color=RED, bold=True)

zones = [
    ("STOP ZONE", "< 0.5m", "즉시 정지", RED),
    ("SLOW ZONE", "< 2m", "0.5m/s 이하 감속", YELLOW),
    ("WARNING ZONE", "< 5m", "경고 + 점진 감속", ACCENT_ORANGE),
]

for i, (name, dist, action, color) in enumerate(zones):
    y = Inches(5.0) + i * Inches(0.65)
    add_shape(slide, Inches(1.3), y, Inches(2.5), Inches(0.5), fill_color=color)
    add_text(slide, Inches(1.5), y + Inches(0.1), Inches(2.1), Inches(0.3),
             f"{name} ({dist})", font_size=13, color=WHITE, bold=True)
    add_text(slide, Inches(4.2), y + Inches(0.1), Inches(8), Inches(0.3),
             f"→ {action}", font_size=14, color=LIGHT_GRAY)

add_text(slide, Inches(1.3), Inches(6.9), Inches(11), Inches(0.3),
         "✅ 비상 정지 거리: 최대속도(1.5m/s)에서 ≤ 1m  |  ISO 3691-4 준수",
         font_size=14, color=GREEN, bold=True)


# ═══════════════════════════════════════════════════════
# SLIDE 13: 선행 확보 + 로드맵 개요
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
         "09  개발 로드맵", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
         "6개월 프로토타입 개발 계획", font_size=30, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.5), Inches(2))

# Accelerators
add_text(slide, Inches(0.8), Inches(1.7), Inches(5), Inches(0.4),
         "✅ 선행 확보 사항 (기존 24개월 → 6개월 압축)", font_size=16, color=GREEN, bold=True)

accelerators = [
    ("자체 제어 시스템", "SW 개발 기간 대폭 단축"),
    ("자율주행 플랫폼", "Navigation Stack 재활용"),
    ("선행 모델 개발 경험", "기계/전기 설계 검증 완료"),
    ("AMR 기본 구동 시스템", "프레임/구동 재설계 최소화"),
]

for i, (title, desc) in enumerate(accelerators):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + col * Inches(6)
    y = Inches(2.2) + row * Inches(0.7)
    add_shape(slide, x, y, Inches(5.7), Inches(0.6), fill_color=CARD_BG)
    add_text(slide, x + Inches(0.2), y + Inches(0.05), Inches(2.5), Inches(0.25),
             f"✅ {title}", font_size=13, color=GREEN, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.3), Inches(5.2), Inches(0.25),
             desc, font_size=12, color=LIGHT_GRAY)

# Timeline - 6 months
add_text(slide, Inches(0.8), Inches(3.8), Inches(5), Inches(0.4),
         "📅 6개월 타임라인", font_size=16, color=ACCENT_ORANGE, bold=True)

phases_compact = [
    ("P1", "설계+소싱", "W1~4", "1개월", ACCENT_BLUE),
    ("P2", "HW통합", "W2~8", "2개월", ACCENT_ORANGE),
    ("P3", "SW통합", "W4~16", "4개월", GREEN),
    ("P4", "테스트", "W12~22", "5.5개월", YELLOW),
    ("P5", "데모", "W20~26", "6개월", RED),
]

# Timeline bars
bar_y = Inches(4.3)
bar_total_width = Inches(11)
total_weeks = 26

for phase, name, period, duration, color in phases_compact:
    # Parse week range
    if "~" in period:
        parts = period.replace("W", "").split("~")
        start_w = int(parts[0])
        end_w = int(parts[1])
    else:
        start_w = int(period.replace("W", ""))
        end_w = start_w
    
    bar_left = Inches(1.5) + (start_w / total_weeks) * bar_total_width
    bar_width = ((end_w - start_w) / total_weeks) * bar_total_width
    if bar_width < Inches(0.5):
        bar_width = Inches(0.5)
    
    add_shape(slide, bar_left, bar_y, bar_width, Inches(0.35), fill_color=color)
    add_text(slide, bar_left + Inches(0.05), bar_y + Inches(0.02), bar_width - Inches(0.1), Inches(0.3),
             f"{phase}", font_size=10, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Week labels
add_text(slide, Inches(1.5), Inches(4.7), Inches(11), Inches(0.25),
         "W1          W4          W8          W12         W16         W20         W26",
         font_size=10, color=MID_GRAY, font_name="Consolas")

# Phase details
add_text(slide, Inches(0.8), Inches(5.1), Inches(5), Inches(0.3),
         "📋 단계별 주요 작업", font_size=14, color=ACCENT_BLUE, bold=True)

phase_details = [
    ("P1 설계+소싱", "양팔 마운트 설계, 섀시 보강, 부품 발주", ACCENT_BLUE),
    ("P2 HW통합", "AMR 보강, 양팔 마운트, 센서/BMS 장착", ACCENT_ORANGE),
    ("P3 SW통합", "양팔 드라이버, MoveIt 2, 협동작업, Fleet 연동", GREEN),
    ("P4 테스트", "성능/안전/통합/하중 검증, 파일럿", YELLOW),
    ("P5 데모", "데모 제작, 인증 준비, 고객 피드백", RED),
]

for i, (title, desc, color) in enumerate(phase_details):
    y = Inches(5.5) + i * Inches(0.38)
    add_text(slide, Inches(1), y, Inches(2.5), Inches(0.3),
             title, font_size=12, color=color, bold=True)
    add_text(slide, Inches(3.8), y, Inches(8.5), Inches(0.3),
             desc, font_size=12, color=LIGHT_GRAY)

# Key assumption
add_shape(slide, Inches(0.8), Inches(7.0), Inches(11.5), Inches(0.3), fill_color=ACCENT_BLUE)
add_text(slide, Inches(1), Inches(7.0), Inches(11), Inches(0.3),
         "핵심: 양팔 COTS(기성품) 선택 → 납기 4~6주  |  SW는 기존 플랫폼 재활용  |  HW/SW 병렬 진행",
         font_size=12, color=WHITE, bold=True)


# ═══════════════════════════════════════════════════════
# SLIDE 14: 리스크 분석
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
         "10  리스크 분석", font_size=14, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(10), Inches(0.6),
         "리스크 매트릭스", font_size=30, color=WHITE, bold=True)
add_divider(slide, Inches(0.8), Inches(1.5), Inches(2))

# Header
risk_headers = ["리스크", "확률", "영향", "대응 방안"]
risk_widths = [Inches(3), Inches(1), Inches(1), Inches(6.5)]

add_shape(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.45), fill_color=ACCENT_BLUE)
x = Inches(0.8)
for h, w in zip(risk_headers, risk_widths):
    add_text(slide, x + Inches(0.1), Inches(1.85), w, Inches(0.35),
             h, font_size=13, color=WHITE, bold=True)
    x += w

risks = [
    ("600kg 하중 안전성 미달", "중", "높음", "Safety PLC 2채널, 조기 안전 테스트"),
    ("양팔 충돌/간섭", "중", "중", "실시간 자기충돌검지, 시뮬레이션 선행"),
    ("배터리 수명 부족", "중", "중", "LiFePO4, 급속 충전, 배터리 스왑"),
    ("AMR 정밀도 미달", "저", "높음", "UWB + LiDAR 이중 위치 추정"),
    ("핵심 부품 수급 불안", "중", "높음", "이중 소싱, 대체 부품 선정"),
    ("안전 인증 지연", "중", "높음", "인증 컨설턴트 조기 투입"),
    ("ROKAE 선제 진입", "저", "높음", "개발 속도 확보, 특허 선점"),
    ("SW 안정성", "중", "중", "CI/CD, HIL 시뮬레이션, 지속적 테스트"),
]

prob_colors = {"저": GREEN, "중": YELLOW, "높음": RED}

for i, (risk, prob, impact, response) in enumerate(risks):
    y = Inches(2.3) + i * Inches(0.58)
    bg = CARD_BG if i % 2 == 0 else DARK_BG
    add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(0.55), fill_color=bg)
    x = Inches(0.8)
    cells = [risk, prob, impact, response]
    widths = risk_widths
    for cell, w in zip(cells, widths):
        color = LIGHT_GRAY
        if cell in prob_colors:
            color = prob_colors[cell]
        add_text(slide, x + Inches(0.1), y + Inches(0.1), w, Inches(0.35),
                 cell, font_size=12, color=color)
        x += w


# ═══════════════════════════════════════════════════════
# SLIDE 15: 마무리
# ═══════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, SECTION_BG)

add_shape(slide, Inches(1.5), Inches(2.2), Inches(2), Pt(4), fill_color=ACCENT_BLUE)
add_text(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1),
         "AMMR-H600", font_size=44, color=WHITE, bold=True)
add_text(slide, Inches(1.5), Inches(3.5), Inches(10), Inches(0.5),
         "양팔 + 30kg+ + 600kg AMR + 자율주행", font_size=22, color=ACCENT_BLUE, bold=True)
add_text(slide, Inches(1.5), Inches(4.2), Inches(10), Inches(0.5),
         "시장에 없는 제품을 만드는 것. 그것이 우리의 목표다.", font_size=18, color=LIGHT_GRAY)

add_text(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.4),
         "v0.2 Draft  |  2026-04-01  |  6개월 프로토타입 개발 목표",
         font_size=14, color=MID_GRAY)


# ═══════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════
output_path = "/root/.openclaw/workspace/AMMR_H600_통합설계서.pptx"
prs.save(output_path)
print(f"✅ PPT saved: {output_path}")
print(f"   Slides: {len(prs.slides)}")
