#!/usr/bin/env python3
"""
에이스엔지니어링 - HMC-SKID 연구개발계획서
20분 → 10분 발표용 PPT 생성 (16장 → 12장 압축)
평가항목 기반: 연구개발내용 우수성(50) + 연구수행능력(25) + 수행계획(15) + 사업화(10)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ===== COLOR PALETTE =====
DARK_NAVY = RGBColor(0x0B, 0x1D, 0x3A)
ACCENT_BLUE = RGBColor(0x00, 0x7B, 0xFF)
ACCENT_CYAN = RGBColor(0x00, 0xBC, 0xD4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x2D, 0x2D, 0x2D)
MID_GRAY = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_BG = RGBColor(0xF4, 0xF7, 0xFA)
CARD_BG = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x00, 0xC8, 0x53)
ORANGE = RGBColor(0xFF, 0x91, 0x00)
RED_ACCENT = RGBColor(0xFF, 0x52, 0x52)
GOLD = RGBColor(0xFF, 0xD6, 0x00)

def add_bg(slide, color=LIGHT_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline(slide, left, top, width, height, lines, font_size=13, color=DARK_GRAY, line_spacing=1.3, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(3)
    return txBox

def add_card(slide, left, top, width, height, title, items, accent_color=ACCENT_BLUE, title_size=13, item_size=11):
    # Card shadow/bg
    add_rounded_rect(slide, left, top, width, height, CARD_BG)
    # Accent top bar
    add_shape(slide, left + Inches(0.15), top + Inches(0.08), Inches(0.06), Inches(0.35), accent_color)
    # Title
    add_textbox(slide, left + Inches(0.35), top + Inches(0.05), width - Inches(0.5), Inches(0.4),
        title, font_size=title_size, color=DARK_NAVY, bold=True)
    # Items
    add_multiline(slide, left + Inches(0.25), top + Inches(0.5), width - Inches(0.4), height - Inches(0.6),
        items, font_size=item_size, color=DARK_GRAY)

def add_top_bar(slide):
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.05), ACCENT_BLUE)

def add_section(slide, title, subtitle="", num=""):
    add_top_bar(slide)
    # Number badge
    if num:
        add_rounded_rect(slide, Inches(0.5), Inches(0.35), Inches(0.65), Inches(0.55), ACCENT_BLUE)
        add_textbox(slide, Inches(0.5), Inches(0.38), Inches(0.65), Inches(0.5),
            num, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, Inches(1.35), Inches(0.35), Inches(10), Inches(0.55),
            title, font_size=26, color=DARK_NAVY, bold=True)
    else:
        add_textbox(slide, Inches(0.5), Inches(0.35), Inches(10), Inches(0.55),
            title, font_size=26, color=DARK_NAVY, bold=True)
    if subtitle:
        add_textbox(slide, Inches(1.35), Inches(0.9), Inches(10), Inches(0.35),
            subtitle, font_size=12, color=MID_GRAY)

def add_page_num(slide, num):
    add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
        f"{num} / 12", font_size=9, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)

def add_eval_tag(slide, text, left=Inches(10.5), top=Inches(0.35)):
    add_rounded_rect(slide, left, top, Inches(2.5), Inches(0.35), RGBColor(0xE3, 0xF2, 0xFD))
    add_textbox(slide, left + Inches(0.1), top + Inches(0.02), Inches(2.3), Inches(0.3),
        text, font_size=9, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ===========================
# SLIDE 1: 표지
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)

# Geometric accents
add_shape(slide, Inches(0), Inches(0), Inches(0.15), prs.slide_height, ACCENT_BLUE)
add_shape(slide, Inches(0), Inches(3.4), prs.slide_width, Inches(0.03), ACCENT_CYAN)

# Top label
add_textbox(slide, Inches(1.5), Inches(1.2), Inches(10), Inches(0.5),
    "글로벌우수기업연구소 육성사업 (GATC)  |  연구개발계획서", font_size=15, color=ACCENT_CYAN)

# Main title
add_textbox(slide, Inches(1.5), Inches(1.9), Inches(10), Inches(1.0),
    "HMC-SKID 기반", font_size=42, color=WHITE, bold=True)
add_textbox(slide, Inches(1.5), Inches(2.8), Inches(10), Inches(1.0),
    "모듈형 컨테이너 솔루션 개발", font_size=42, color=WHITE, bold=True)

add_textbox(slide, Inches(1.5), Inches(3.7), Inches(10), Inches(0.5),
    "Development of HMC SKID-based Modular Containerized Solution", font_size=16, color=RGBColor(0x88, 0xAA, 0xDD))

# Info bar
add_shape(slide, Inches(1.5), Inches(5.0), Inches(10), Inches(0.02), ACCENT_BLUE)
info_items = [
    "주관: ㈜에이스엔지니어링  |  연구책임: 주재흥 이사",
    "사업기간: 2026 ~ 2028 (3개년)  |  총 R&D 투자: 45.1억원",
]
add_multiline(slide, Inches(1.5), Inches(5.2), Inches(10), Inches(0.8), info_items,
    font_size=14, color=RGBColor(0x99, 0xBB, 0xDD))


# ===========================
# SLIDE 2: 현장 문제 & 기술개발 필요성 (A. 기술개발 적정성 20점)
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_section(slide, "현장 문제점 & 기술개발 필요성", "시장이 직면한 기술적 격차와 HMC-SKID의 해결 방향", "01")
add_eval_tag(slide, "A. 기술개발의 적정성 (20점)")

# Problem cards
problems = [
    ("복잡한 현장 작업", RED_ACCENT, [
        "• Containerized ESS의 복잡한 현장 배관·전기 작업",
        "• 설치 기간 장기화 (평균 2~3개월 소요)",
        "• 높은 현장 노무비 발생",
    ]),
    ("품질 편차 위험", ORANGE, [
        "• 현장 작업자 숙련도에 따른 품질 불일치",
        "• 재작업률 증가 → 비용·일정 초과",
        "• 안전 사고 위험 상존",
    ]),
    ("표준화 부재", DARK_NAVY, [
        "• 대규모 ESS 배포용 표준 사전제작 SKID 플랫폼 없음",
        "• 프로젝트마다 개별 현장 엔지니어링 필요",
        "• 시장에 Plug-and-Play 솔루션 부재",
    ]),
]

for i, (title, color, items) in enumerate(problems):
    x = Inches(0.5) + Inches(i * 4.2)
    add_card(slide, x, Inches(1.7), Inches(3.9), Inches(2.3), title, items, accent_color=color)

# Solution highlight
add_rounded_rect(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(2.8), CARD_BG)
add_shape(slide, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.06), ACCENT_BLUE)

add_textbox(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.5),
    "해결책: HMC-SKID (Hybrid Modular Connector - SKID) 기술", font_size=20, color=ACCENT_BLUE, bold=True)

solutions = [
    "Plug-and-Play 방식 배포 → 현장 연결 작업 획기적 제거",
    "사전제작형 표준화 SKID 플랫폼 → 설치 기간 30% 단축 (60일 → 42일)",
    "품질 표준화로 현장 결함 억제 (40% 개선) + 노무비 25~35% 절감",
    "글로벌 ESS 시장 CAGR 60%+ 초고성장 시장 대응 필수 기술",
]
add_multiline(slide, Inches(1.0), Inches(5.1), Inches(11.0), Inches(1.8), solutions,
    font_size=13, color=DARK_GRAY)

add_page_num(slide, 2)


# ===========================
# SLIDE 3: 핵심기술 & 개발목표 (A. 연구개발내용 구체성 15점)
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_section(slide, "핵심기술 & 개발 목표", "3대 핵심 기술 영역 및 연구 목표", "02")
add_eval_tag(slide, "A. 연구개발내용의 구체성 (15점)")

techs = [
    ("HMC-SKID 모듈형 솔루션", ACCENT_BLUE, [
        "• 자동 정렬 보정 기능 포함 자동화 기계 결합",
        "• 누수 제로 연결 메커니즘 & 고신뢰성 조립",
        "• 정밀 정렬 제어 시스템 구축",
        "• HV 자동 커넥터 1,000회+ 사이클 테스트",
    ]),
    ("Grid-Forming 인버터 기술", GREEN, [
        "• 동기발전기 유사 그리드 지원 기능 (GFM)",
        "• 동기발전기 특성 에뮬레이션 핵심 알고리즘",
        "• 전압→주파수→전력흐름 단계별 제어",
        "• MATLAB 모델링 → H/W-in-the-Loop 테스트",
    ]),
    ("AI 데이터센터 응용 시제품", ORANGE, [
        "• 열 프로필 기반 열/구조 설계",
        "• GFM 기능 통합 명세 확립",
        "• 유럽형 (2차년도) + 미국형 (3차년도) 시제품",
        "• 대규모 재생에너지·데이터센터 타겟",
    ]),
]

for i, (title, color, items) in enumerate(techs):
    x = Inches(0.5) + Inches(i * 4.2)
    add_card(slide, x, Inches(1.7), Inches(3.9), Inches(3.0), title, items, accent_color=color)

# Expected outcomes
add_card(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(2.0), "연도별 기대 결과물", [
    "1차년도 (2026): GFM 인버터 기초 설계 완료 + HMC 프로토타입 메커니즘 검증 → 2차년도 HW 구현 기반 마련",
    "2차년도 (2027): HMC-SKID 시스템 제작·성능검증 + GFM 알고리즘 정정 + 유럽형 사양 컨테이너 시제품 제작",
    "3차년도 (2028): GFM 인버터 통합 HMC-SKID 모듈형 Containerized Solution 완성 → 글로벌 배포 기술 준비 완료",
], accent_color=DARK_NAVY, item_size=11)

add_page_num(slide, 3)


# ===========================
# SLIDE 4: 개발 전략 & 위험관리 (C. 수행전략·방법 10점)
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_section(slide, "개발 전략 & 위험 관리", "통합 개발 접근법: 설계→시제품→검증→최적화", "03")
add_eval_tag(slide, "C. 연구개발 수행 전략 및 방법 (10점)")

# Development path
add_card(slide, Inches(0.5), Inches(1.7), Inches(6.0), Inches(2.5), "HMC-SKID 개발 경로", [
    "① 설계 최적화 → ② 시제품 제작 → ③ 성능 평가 → ④ 반복 개선",
    "",
    "위험 완화: 1,000회 이상 자동 정렬 메커니즘 신뢰성 테스트",
    "모듈형 표준화 설계 적용으로 재현성·확장성 확보",
    "협력: EPC Power 및 연구기관과 정기 통합 기술 미팅",
], accent_color=ACCENT_BLUE)

add_card(slide, Inches(6.8), Inches(1.7), Inches(6.0), Inches(2.5), "Grid-Forming 개발 경로", [
    "① 알고리즘 설계 → ② MATLAB 모델링 → ③ H/W-in-the-Loop 테스트",
    "",
    "전압 제어 → 주파수 제어 → 전력흐름 제어 단계별 검증",
    "상용화 전 MATLAB 시뮬레이션으로 안정성 확보",
    "제한적 현장 테스트를 통한 최종 검증",
], accent_color=GREEN)

# Risk management
add_card(slide, Inches(0.5), Inches(4.5), Inches(6.0), Inches(2.5), "기술적 위험 완화", [
    "• HV 자동 커넥터 신뢰성: 중복 안전 메커니즘 설계",
    "• GFM 제어 안정성: 시뮬레이션 + HIL 테스트 + 현장 실험",
    "• 국제 인증: UL·CE 인증기관과 조기 협력",
], accent_color=RED_ACCENT)

add_card(slide, Inches(6.8), Inches(4.5), Inches(6.0), Inches(2.5), "자원·일정 위험 완화", [
    "• 핵심 인력: 백업 리더 교차 교육 & 지식 공유",
    "• 공급망: 핵심 부품 다원 공급처 + 전략적 재고",
    "• 일정: 마일스톤 대비 월간 검토 + 분기별 위험 재평가",
], accent_color=ORANGE)

add_page_num(slide, 4)


# ===========================
# SLIDE 5: 실행 일정 (C. 추진일정 구체성 5점)
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_section(slide, "실행 일정 (Execution Schedule)", "명확한 마일스톤 기반 구조화된 개발 일정", "04")
add_eval_tag(slide, "C. 추진일정의 구체성 (5점)")

# Timeline
phases = [
    ("1차년도 (2026)", ACCENT_BLUE, [
        "Q2: HMC-SKID 설계 + GFM 알고리즘 기초",
        "Q3: 시제품 착수",
        "Q4: 첫 성능 평가 + 설계 반복 개선",
        "★ Gate Review 1",
    ]),
    ("2차년도 (2027)", GREEN, [
        "Q1~Q2: 전체 서브시스템 기능·안전성 테스트",
        "Q3: 유럽형 사양 컨테이너 시제품 제작",
        "Q4: 통합 시스템 테스트·검증",
        "★ EU 시제품 완성",
    ]),
    ("3차년도 (2028)", ORANGE, [
        "Q1~Q2: 미국형 사양 시제품 완성·통합",
        "Q3: 공장 인수 테스트 (FAT)·시스템 검증",
        "Q4: 최종 결과물 + 기술이전 착수",
        "★ 상용화 준비 완료",
    ]),
]

for i, (title, color, items) in enumerate(phases):
    x = Inches(0.5) + Inches(i * 4.2)
    add_card(slide, x, Inches(1.7), Inches(3.9), Inches(3.0), title, items, accent_color=color)

# Timeline bar
add_shape(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(0.12), RGBColor(0xE0, 0xE0, 0xE0))
add_shape(slide, Inches(0.8), Inches(5.2), Inches(3.8), Inches(0.12), ACCENT_BLUE)
add_shape(slide, Inches(4.6), Inches(5.2), Inches(3.9), Inches(0.12), GREEN)
add_shape(slide, Inches(8.5), Inches(5.2), Inches(4.0), Inches(0.12), ORANGE)

milestones = [
    (Inches(0.8), "설계+기초", ACCENT_BLUE),
    (Inches(2.5), "시제품착수", ACCENT_BLUE),
    (Inches(4.6), "서브시스템테스트", GREEN),
    (Inches(6.5), "EU시제품", GREEN),
    (Inches(8.5), "미국시제품", ORANGE),
    (Inches(10.5), "FAT+기술이전", ORANGE),
]
for x, label, color in milestones:
    add_shape(slide, x, Inches(5.1), Inches(0.12), Inches(0.3), color)
    add_textbox(slide, x - Inches(0.3), Inches(5.45), Inches(0.8), Inches(0.4),
        label, font_size=8, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Gantt-like bottom
add_card(slide, Inches(0.5), Inches(5.9), Inches(12.3), Inches(1.1), "핵심 마일스톤", [
    "Gate Review 1 (2026 Q4) → EU 시제품 완성 (2027 Q4) → FAT 통과 (2028 Q3) → 기술이전 착수 (2028 Q4) → 상용화 (2029+)",
], accent_color=DARK_NAVY, item_size=11)

add_page_num(slide, 5)


# ===========================
# SLIDE 6: 연구팀 & 수행능력 (B. 연구인력 전문성 15점)
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_section(slide, "연구팀 구성 & 수행 능력", "글로벌 사업화 경험 보유 전문가 조직", "05")
add_eval_tag(slide, "B. 연구인력의 전문성 (15점)")

# Lead researcher
add_rounded_rect(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(1.3), CARD_BG)
add_shape(slide, Inches(0.5), Inches(1.7), Inches(0.08), Inches(1.3), ACCENT_BLUE)
add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.4),
    "책임연구자: 주재흥 이사 (주관연구개발기관)", font_size=16, color=DARK_NAVY, bold=True)
add_multiline(slide, Inches(0.8), Inches(2.25), Inches(11.5), Inches(0.6), [
    "ESS 시스템 20년 이상 경력의 컨테이너 기술 전문가 | 구조설계 총괄 17년 경력 | 프로젝트 기술적 방향 설정 및 국제 협력 총괄",
], font_size=12, color=MID_GRAY)

# Core team
add_card(slide, Inches(0.5), Inches(3.3), Inches(4.0), Inches(3.7), "핵심 연구팀", [
    "중앙기술연구소: 전체 23명",
    "연구인력: 9명 (석사 2, 학사 7)",
    "평균연구경력: 10년 이상",
    "",
    "• 조용래 차장 — 전기설계 (26년)",
    "• 양수석 차장 — 전기설계 (26년)",
    "• 장희수 차장 — 전기설계 (14년)",
    "• 마상진 과장 — 구조설계 (13년)",
    "• 성동근 과장 — Materials Science",
    "  (석사, 독일 Bochum)",
], accent_color=ACCENT_BLUE)

# Partner orgs
add_card(slide, Inches(4.8), Inches(3.3), Inches(4.0), Inches(3.7), "공동연구기관", [
    "EPC Power — GFM 인버터 공동개발",
    "",
    "R&D 자문·시험평가:",
    "  한국전기연구원 (KERI)",
    "  한국생산기술연구원 (KITECH)",
    "",
    "산학협력:",
    "  서울대학교 — GFM 핵심기술",
    "  성균관대학교 — 전력시스템",
    "  동아대학교 — 인재육성",
], accent_color=GREEN)

# Supporting
add_card(slide, Inches(9.1), Inches(3.3), Inches(3.7), Inches(3.7), "지원체계", [
    "지식재산 전담:",
    "  이호용 상무 (전략TF장)",
    "  前 한화큐셀 특허전략 경험",
    "  파이특허법률사무소 (10년+)",
    "",
    "인증 지원:",
    "  UL·CE 조기 협력",
    "  ISO 9001/14001/45001",
    "",
    "Fluence — 글로벌 실증·사업화",
], accent_color=ORANGE)

add_page_num(slide, 6)


# ===========================
# SLIDE 7: 연구시설 & 인프라 (B. 시설·장비 10점)
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_section(slide, "연구시설 & 인프라", "최첨단 개발 인프라 보유 및 확장 계획", "06")
add_eval_tag(slide, "B. 연구시설 및 장비 보유현황 (10점)")

facilities = [
    ("기존 연구시설", ACCENT_BLUE, [
        "CSC 테스트 시설: 16.2m × 5.9m × 5.6m",
        "종합 구조·하중 테스트 가능",
        "",
        "기계 설계 워크스테이션",
        "ANSYS, CATIA, MIDAS NFX",
        "",
        "전기 테스트 랩",
        "고전압·전력전자 성능 검증",
        "",
        "열 관리 테스트 챔버",
        "-30°C ~ +70°C 극한 환경 테스트",
    ]),
    ("계획된 투자 (2027)", GREEN, [
        "총 약 30억원 규모 신규 인프라",
        "",
        "배터리 충방전기 (1000V/300A)",
        "배터리팩 챔버",
        "신속 진단기 (SOH/SOC/SOP)",
        "배터리팩 검사시스템",
        "자동 적재 시스템 (연 100EA)",
        "배터리팩 칠러 (-30~70°C)",
    ]),
    ("협력 접근 시설", ORANGE, [
        "KERI: 그리드 표준 준수 테스트",
        "KITECH: 제조 공정 최적화",
        "SNU/SKKU: 고급 시뮬레이션",
        "",
        "자동 로딩 시스템 (10억원)",
        "신속한 시제품 테스트 사이클",
        "",
        "협력기관 보유 장비 활용으로",
        "자체 투자 대비 효율적 운영",
    ]),
]

for i, (title, color, items) in enumerate(facilities):
    x = Inches(0.5) + Inches(i * 4.2)
    add_card(slide, x, Inches(1.7), Inches(3.9), Inches(5.3), title, items, accent_color=color, item_size=11)

add_page_num(slide, 7)


# ===========================
# SLIDE 8: 시장 기회 & 경쟁우위 (A. 사업화 가능성 15점)
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_section(slide, "시장 기회 & 경쟁 우위", "글로벌 ESS 초고성장 시장에서의 차별화 포지셔닝", "07")
add_eval_tag(slide, "A. 성과활용 및 사업화 가능성 (15점)")

# Market data
add_card(slide, Inches(0.5), Inches(1.7), Inches(6.0), Inches(2.5), "시장 규모 & 성장성", [
    "글로벌 ESS 시장: CAGR 60% 이상 초고성장",
    "AI Data Center 성장: CAGR 30%+",
    "2024년 글로벌 ESS 시장 → 2030년 50B+ USD 전망",
    "",
    "핵심: 설치 시간 단축·노무비 절감·품질 개선이 시장 요구사항",
], accent_color=ACCENT_BLUE)

# Competitive advantage
add_card(slide, Inches(6.8), Inches(1.7), Inches(6.0), Inches(2.5), "경쟁 우위 (HMC 기술 독특성)", [
    "현장 연결 제거 유일한 사전제작형 SKID 솔루션",
    "Grid-Forming 고급 제어 기술로 그리드 안정성 강화",
    "설계→배포 수직 통합으로 운영 효율·품질 극대화",
    "배포 속도 혁신 → 시장 선도 가능",
], accent_color=GREEN)

# Cost-benefit
add_card(slide, Inches(0.5), Inches(4.5), Inches(6.0), Inches(2.5), "비용-편익 분석", [
    "설치 시간: 30% 단축 (60일 → 42일)",
    "노무비: 25~35% 감소",
    "품질 개선: 현장 결함 40% 감소",
    "총소유비용(TCO): 15~20% 절감",
], accent_color=ORANGE)

# Target market
add_card(slide, Inches(6.8), Inches(4.5), Inches(6.0), Inches(2.5), "타겟 시장 & 진출 계획", [
    "2027: 유럽 시장 진출 (Hitachi 유럽 협력)",
    "2027: 북미 상용 출시 (EPC Power 채널)",
    "2028: APAC 시장 확대",
    "대상: 대규모 재생에너지·AI 데이터센터·그리드 현대화",
], accent_color=DARK_NAVY)

add_page_num(slide, 8)


# ===========================
# SLIDE 9: 투자 & 수익 전망
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_section(slide, "투자 요약 & 수익 전망", "보수적 가정 기반 명확한 수익 실현", "08")
add_eval_tag(slide, "A. 성과활용 및 사업화 가능성 (15점)")

# Investment summary
add_card(slide, Inches(0.5), Inches(1.7), Inches(6.0), Inches(2.5), "투자 요약", [
    "총 R&D 투자: 45.1억원",
    "  정부 지원금: 33억원 (73%)",
    "  ACE Engineering: 12.1억원 (27%)",
    "",
    "2027년 신규 연구장비 구축: 약 30억원 별도",
], accent_color=ACCENT_BLUE)

# Revenue projection
add_card(slide, Inches(6.8), Inches(1.7), Inches(6.0), Inches(2.5), "예상 수익 (Revenue Impact)", [
    "5년 누적 수익: 600~1,000억원",
    "10년 NPV: 1,500억원+",
    "직접 일자리 창출: 50~100명",
    "공급망 파급효과: 200억원+",
    "",
    "2028년: 10~15 unit 판매 목표",
    "2029~2030년: 50~80 unit 판매 목표",
], accent_color=GREEN)

# KPI table
kpi_data = [
    ("항목", "2026", "2027", "2028", "2029", "2030"),
    ("판매(unit)", "-", "초기", "10~15", "50~80", "80+"),
    ("매출(억원)", "-", "PoC", "50~100", "300~500", "500+"),
    ("특허(건)", "2", "3", "3", "4", "4"),
    ("연구인력", "12", "15", "18", "21", "25"),
]

col_widths = [Inches(2.0), Inches(1.6), Inches(1.6), Inches(1.6), Inches(1.6), Inches(1.6)]
x_start = Inches(1.7)
y_start = Inches(4.7)

for row_idx, row in enumerate(kpi_data):
    x = x_start
    y = y_start + Inches(row_idx * 0.48)
    for col_idx, cell in enumerate(row):
        w = col_widths[col_idx]
        if row_idx == 0:
            bg = DARK_NAVY
            txt_color = WHITE
            is_bold = True
        else:
            bg = CARD_BG if row_idx % 2 == 1 else RGBColor(0xEE, 0xF2, 0xF7)
            txt_color = DARK_GRAY
            is_bold = col_idx == 0
        add_shape(slide, x, y, w, Inches(0.44), bg)
        add_textbox(slide, x, y + Inches(0.05), w, Inches(0.35), cell,
            font_size=11, color=txt_color, bold=is_bold, alignment=PP_ALIGN.CENTER)
        x += w

add_page_num(slide, 9)


# ===========================
# SLIDE 10: 사업화 전략 (D. 성과활용·사업화 계획 10점)
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_section(slide, "기술전환 & 사업화 전략", "Post Project — 시장 진입 가속화 단계별 전략", "09")
add_eval_tag(slide, "D. 성과활용·사업화 계획 (10점)")

# Tech transfer
add_card(slide, Inches(0.5), Inches(1.7), Inches(4.0), Inches(2.5), "기술 전환 로드맵", [
    "① 개념증명 (PoC)",
    "   HMC-SKID + GFM 알고리즘 검증",
    "",
    "② 기술 실증",
    "   최종사용자 평가용 동작 시제품",
    "",
    "③ 상용화 준비",
    "   대량 생산 준비 완료 시스템",
], accent_color=ACCENT_BLUE)

# 3-stage commercialization
stages = [
    ("1단계: 기반·검증", GREEN, [
        "• EPC Power 고객 2~3개 프로젝트",
        "• UL 9540, IEC 62619 인증",
        "• 기술 문서·교육 프로그램 개발",
    ]),
    ("2단계: 확장·파트너십", ORANGE, [
        "• 유럽 SI 유통망 구축",
        "• 북미 상용 출시",
        "• APAC 시장 개발 착수",
    ]),
    ("3단계: 스케일·생태계", DARK_NAVY, [
        "• 연 50건+ 판매 달성",
        "• 3~5개 국제 SI 라이센싱",
        "• 설계 컨설팅 서비스 창출",
    ]),
]

for i, (title, color, items) in enumerate(stages):
    x = Inches(4.8) + Inches(i * 2.9)
    add_card(slide, x, Inches(1.7), Inches(2.7), Inches(2.5), title, items, accent_color=color, item_size=10)

# IP protection
add_card(slide, Inches(0.5), Inches(4.5), Inches(6.0), Inches(2.5), "IP 보호 전략", [
    "국내 특허 출원/등록 → 국제 특허 (PCT) 출원",
    "핵심 기술 강력한 특허 포트폴리오 구축",
    "기술 영업비밀 등록으로 글로벌 경쟁 우위 확보",
    "Grid-Forming 인버터: 태양광·재생에너지 시스템 편익",
    "FTO 조사 + NDA + 접근권한 제한으로 기술유출 방지",
], accent_color=RED_ACCENT)

# Success drivers
add_card(slide, Inches(6.8), Inches(4.5), Inches(6.0), Inches(2.5), "성공 핵심 요인", [
    "초기 기준 성공: EPC Power 프로젝트 성공 → 시장 신뢰·채택 드라이브",
    "파트너 생태계: 글로벌 유통·설치 파트너 개발 → 시장 침투 가속화",
    "지속적 개선: 현장 피드백 기반 제품 성능 최적화 프로세스 운영",
    "→ 기술독특성 × 시장타이밍 × 파트너십 = 글로벌 선도",
], accent_color=ACCENT_CYAN)

add_page_num(slide, 10)


# ===========================
# SLIDE 11: IP & 인증 전략
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_BG)
add_section(slide, "지식재산 & 인증 현황", "특허 포트폴리오 강화 및 국제 인증 대응", "10")

# IP portfolio
add_card(slide, Inches(0.5), Inches(1.7), Inches(6.0), Inches(2.8), "지식재산권 현황", [
    "특허등록: 11건",
    "특허출원: 2건",
    "해외특허 (미국·중국·PCT): 3건",
    "",
    "전략TF: 이호용 상무 (前 한화큐셀 특허전략)",
    "선행특허 조사: 파이특허법률사무소 (10년+ 협업)",
    "사내공모전 → 신규 아이템 발굴 → IP 출원",
], accent_color=ACCENT_BLUE)

# Certifications
add_card(slide, Inches(6.8), Inches(1.7), Inches(6.0), Inches(2.8), "인증 현황 & 계획", [
    "보유 인증:",
    "  ISO 9001 / 14001 / 45001 / 27001",
    "  한국선급 (KR) · LRQA · IMO · ABS",
    "",
    "목표 인증:",
    "  UL 9540 (미국 안전인증)",
    "  IEC 62619 (배터리 안전)",
    "  CE (유럽 적합성)",
], accent_color=GREEN)

# National R&D
add_card(slide, Inches(0.5), Inches(4.8), Inches(12.3), Inches(2.2), "국가연구개발사업 수행실적", [
    "ESS 설치공간 화재예방·차단 시스템개발 ('21~'25) — 산업통상자원부",
    "MWh급 선박용 고안전성 LiB-ESS 통합시스템 국산화 ('21~'23) — 해양수산부",
    "야전병원 ICT융합 플랫폼 디자인개발 ('21~'25) — 국방부",
    "",
    "→ 3건의 국가R&D 수행 경험으로 GATC 과제 수행 역량 입증",
], accent_color=DARK_NAVY, item_size=12)

add_page_num(slide, 11)


# ===========================
# SLIDE 12: 핵심요약 & 마무리
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_NAVY)

add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.05), ACCENT_CYAN)

add_textbox(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.7),
    "핵심 요약", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

summaries = [
    ("기술 독특성", ACCENT_BLUE, [
        "현장 연결 제거 유일한 SKID 솔루션",
        "Grid-Forming 고급 제어 기술",
        "Plug-and-Play 배포 혁신",
        "설치시간 30%·노무비 30%·TCO 18% 절감",
    ]),
    ("시장 기회", GREEN, [
        "글로벌 ESS 시장 CAGR 60%+",
        "AI 데이터센터 CAGR 30%+",
        "5년 누적 수익 600~1,000억원",
        "10년 NPV 1,500억원+",
    ]),
    ("수행 역량", ORANGE, [
        "45년 컨테이너 역사 + ESS 글로벌 1위",
        "20년+ ESS 전문가 연구팀",
        "13건 특허 + 국제 인증 체계",
        "국가R&D 3건 수행 경험",
    ]),
]

for i, (title, color, items) in enumerate(summaries):
    x = Inches(0.5) + Inches(i * 4.2)
    # Card with dark bg
    add_rounded_rect(slide, x, Inches(1.6), Inches(3.9), Inches(3.2), RGBColor(0x12, 0x2A, 0x4F))
    # Accent line
    add_shape(slide, x + Inches(0.2), Inches(1.75), Inches(0.06), Inches(0.4), color)
    add_textbox(slide, x + Inches(0.45), Inches(1.75), Inches(3.2), Inches(0.4),
        title, font_size=18, color=color, bold=True)
    add_multiline(slide, x + Inches(0.35), Inches(2.3), Inches(3.2), Inches(2.2),
        items, font_size=12, color=RGBColor(0xCC, 0xDD, 0xEE))

# Bottom CTA
add_shape(slide, Inches(0), Inches(5.4), prs.slide_width, Inches(0.03), ACCENT_CYAN)

add_textbox(slide, Inches(1), Inches(5.7), Inches(11), Inches(0.6),
    "HMC-SKID = 글로벌 ESS 시장의 게임 체인저", font_size=22, color=ACCENT_CYAN, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.6),
    "경청해 주셔서 감사합니다", font_size=20, color=RGBColor(0x99, 0xBB, 0xDD), alignment=PP_ALIGN.CENTER)

add_page_num(slide, 12)


# ===== SAVE =====
output_path = "/root/.openclaw/workspace/HMC-SKID_연구개발계획서_10분.pptx"
prs.save(output_path)
print(f"✅ PPT saved: {output_path}")
print(f"   Slides: {len(prs.slides)}")
