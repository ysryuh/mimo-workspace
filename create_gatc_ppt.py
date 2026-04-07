#!/usr/bin/env python3
"""
에이스엔지니어링 - 글로벌우수기업연구소 육성사업 (GATC)
10분 발표용 PPT 생성 스크립트 (21장 → 12장 압축)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ===== COLOR PALETTE =====
DARK_BLUE = RGBColor(0x00, 0x2B, 0x5C)
ACCENT_BLUE = RGBColor(0x00, 0x6E, 0xB8)
LIGHT_BLUE = RGBColor(0xE8, 0xF4, 0xFD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
GREEN = RGBColor(0x00, 0x96, 0x88)
ORANGE = RGBColor(0xFF, 0x8F, 0x00)

def add_bg(slide, color=DARK_BLUE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=18, color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT, font_name='맑은 고딕'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_textbox(slide, left, top, width, height, items, font_size=14, color=DARK_GRAY, bold_first=False, font_name='맑은 고딕', line_spacing=1.5):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(4)
        if bold_first and '•' not in item and ':' not in item:
            p.font.bold = True
    return txBox

def add_section_header(slide, text, sub_text=""):
    # Top accent bar
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_BLUE)
    # Left accent
    add_shape(slide, Inches(0.6), Inches(1.0), Inches(0.08), Inches(0.5), ACCENT_BLUE)
    add_textbox(slide, Inches(0.9), Inches(0.8), Inches(10), Inches(0.7), text, font_size=28, color=DARK_BLUE, bold=True)
    if sub_text:
        add_textbox(slide, Inches(0.9), Inches(1.5), Inches(10), Inches(0.4), sub_text, font_size=14, color=MID_GRAY)

def add_card(slide, left, top, width, height, title, items, title_color=DARK_BLUE, bg_color=WHITE):
    # Card background
    card = add_shape(slide, left, top, width, height, bg_color)
    # Card shadow effect via border
    card.shadow.inherit = False
    # Title bar
    add_shape(slide, left, top, width, Inches(0.5), title_color)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.05), width - Inches(0.4), Inches(0.45), title, font_size=14, color=WHITE, bold=True)
    # Content
    add_bullet_textbox(slide, left + Inches(0.2), top + Inches(0.6), width - Inches(0.4), height - Inches(0.8), items, font_size=11, color=DARK_GRAY)

def add_page_number(slide, num):
    add_textbox(slide, Inches(12.5), Inches(7.0), Inches(0.6), Inches(0.4), str(num), font_size=10, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)


# ===========================
# SLIDE 1: 표지
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)

# Accent stripes
add_shape(slide, Inches(0), Inches(2.8), prs.slide_width, Inches(0.06), ACCENT_BLUE)
add_shape(slide, Inches(0), Inches(4.6), prs.slide_width, Inches(0.04), ACCENT_BLUE)

# Main title
add_textbox(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(0.6),
    "글로벌우수기업연구소 육성사업 (GATC)", font_size=18, color=RGBColor(0xAA, 0xCC, 0xEE), bold=False)

add_textbox(slide, Inches(1.5), Inches(1.6), Inches(10), Inches(1.2),
    "기업부설연구소 성장전략계획서", font_size=40, color=WHITE, bold=True)

add_textbox(slide, Inches(1.5), Inches(3.1), Inches(10), Inches(0.8),
    "글로벌에너지·전력인프라의 표준이 되는\nModular Containerized Solution", font_size=22, color=RGBColor(0xCC, 0xDD, 0xEE))

# Company info
add_textbox(slide, Inches(1.5), Inches(5.0), Inches(5), Inches(0.4),
    "주관연구개발기관  ㈜에이스엔지니어링  |  중앙기술연구소", font_size=14, color=RGBColor(0x99, 0xBB, 0xDD))
add_textbox(slide, Inches(1.5), Inches(5.5), Inches(5), Inches(0.4),
    "연구소장  주재흥 이사  |  사업기간  2026 ~ 2029", font_size=14, color=RGBColor(0x99, 0xBB, 0xDD))


# ===========================
# SLIDE 2: 목차
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), Inches(4.5), prs.slide_height, DARK_BLUE)

add_textbox(slide, Inches(0.5), Inches(2.5), Inches(3.5), Inches(1.5),
    "AGENDA", font_size=44, color=WHITE, bold=True)

sections = [
    ("01", "기업 개요 및 R&D 역량", "사업화능력, 연구인력, 핵심기술"),
    ("02", "글로벌 성장 로드맵", "비전, 5단계 성장전략, 핵심활동"),
    ("03", "R&D 투자·인력·인프라 계획", "재원조달, 글로벌인력, 장비구축"),
    ("04", "협력 네트워크", "국내외 R&D·비즈니스 네트워크"),
    ("05", "부합성·경영진 의지·핵심요약", "기업-연구소 부합, 성장 KPI"),
]

for i, (num, title, desc) in enumerate(sections):
    y = Inches(1.2) + Inches(i * 1.15)
    # Number circle
    add_textbox(slide, Inches(5.2), y, Inches(0.8), Inches(0.5), num, font_size=28, color=ACCENT_BLUE, bold=True)
    add_textbox(slide, Inches(6.0), y, Inches(6), Inches(0.4), title, font_size=20, color=DARK_BLUE, bold=True)
    add_textbox(slide, Inches(6.0), y + Inches(0.4), Inches(6), Inches(0.35), desc, font_size=12, color=MID_GRAY)

add_page_number(slide, 2)


# ===========================
# SLIDE 3: 기업 개요 & 재무현황 (A. 사업화능력 25점)
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_section_header(slide, "01  기업 개요 & 재무현황", "A. 사업화능력 — 매출 실적 (15점) + R&D 투자 (10점)")

# Left: Company overview card
add_card(slide, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.8), "기업 개요", [
    "• 설립: 1981년 (컨테이너산업 45년 역사)",
    "• 사업: ESS Containerized Solution (Top-tier)",
    "• 매출: 2024년 6,340억원 (연평균 64% 성장)",
    "• 종업원: 111명 (연평균 24% 성장)",
    "• 글로벌: 30개국 이상 제품 공급 경험",
    "• 주요 거래처: Fluence(65%), 효성중공업(34%)",
    "• 미국·베트남 법인 운영",
    "",
    "핵심 인증: ISO 9001/14001/45001/27001",
    "· KR·LRQA · 한국선급 · IMO · ABS",
])

# Right: Financial data card
add_card(slide, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.2), "재무 현황 (2022~2024)", [
    "• 매출액: 2,354억 → 2,959억 → 6,340억 (연평균 64%)",
    "• 총자산: 425억 → 1,930억 → 3,532억 (연평균 22%)",
    "• 영업이익률: 6% → 8% → 8%",
    "• 2025년 매출 목표: 4,500억원",
])

# R&D investment card
add_card(slide, Inches(6.8), Inches(4.6), Inches(5.8), Inches(2.4), "R&D 투자 (A항목 충족)", [
    "• 기존 연 5억 → 연 10억원 수준으로 확대",
    "• 자체조달 50% / 국가R&D 20% / 수요처지원 20%",
    "• 주주증자 550억 확보 (2025년)",
    "• 데이터센터·스키드·GFM인버터 우선투자",
    "• 매출 무관 R&D 연 5억+ 하한선 유지",
])

add_page_number(slide, 3)


# ===========================
# SLIDE 4: 연구소 역량 & 인력 (B. 연구인력 구성 15점)
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_section_header(slide, "01  연구소 역량 & 인력 구성", "B. 연구인력 구성 — 연구인력 (10점) + 연구소장 자격 (5점)")

# Org structure card
add_card(slide, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.8), "중앙기술연구소 조직", [
    "연구소장: 주재흥 이사 (구조설계 총괄, 17년 경력)",
    "",
    "• 연구인력: 9명 (석사 2명, 학사 7명)",
    "• 총 인력규모: 중앙기술연구소 전체 23명",
    "• 전문분야: 구조설계·전기설계·구조해석",
    "• 평균연구경력: 10년 이상",
    "",
    "핵심 인력:",
    "  조용래 차장 — 전기설계 전문 (26년 경력)",
    "  양수석 차장 — 전기설계 전문 (26년 경력)",
    "  장희수 차장 — 전기설계 (14년 경력)",
    "  마상진 과장 — 구조설계 (13년 경력)",
    "  성동근 과장 — Materials Science (석사, 독일 Bochum)",
])

# IP & Cert card
add_card(slide, Inches(6.8), Inches(2.2), Inches(5.8), Inches(2.2), "지식재산권 & 인증 현황", [
    "• 특허등록: 11건 / 특허출원: 2건",
    "• 해외특허 (미국·중국·PCT): 3건",
    "• ISO 9001/14001/45001/27001 인증",
    "• 한국선급·LRQA·IMO·ABS 인증",
])

# IP strategy card
add_card(slide, Inches(6.8), Inches(4.6), Inches(5.8), Inches(2.4), "지식재산 전담 체계", [
    "• 전략TF장: 이호용 상무 (前 한화큐셀 특허전략 경험)",
    "• 사내연구소: IP 대상 아이템 R&D 수행",
    "• 사내공모전 기반 신규 아이템 발굴",
    "• 파이특허법률사무소: 선행특허 조사·출원 (10년+ 협업)",
])

add_page_number(slide, 4)


# ===========================
# SLIDE 5: 핵심기술 & 수상실적 (E. 수상실적 25점)
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_section_header(slide, "01  핵심기술 & 기술사업화 성과", "E. 수상실적 — 특허·인증·연구과제 수행실적 (25점)")

# Tech commercialization
add_card(slide, Inches(0.6), Inches(2.2), Inches(7.6), Inches(4.8), "기술사업화 대표 실적", [
    "Powin Centipede (2021~) — 누적 3,500억원",
    "Fluence Gen6 Cube (2020~2025) — 누적매출 9,000억원",
    "Fluence GSP 2000 — 누적 150억 + 수주잔고 700억",
    "Fluence GSP 5000 — 누적 1,200억 + 수주잔고 3,000억",
    "Fluence Smartstack — 누적 50억 + 수주잔고 300억",
    "",
    "국가연구개발사업 수행실적:",
    "• ESS 설치공간 화재예방·차단 시스템개발 ('21~'25)",
    "• MWh급 선박용 고안전성 LiB-ESS 통합시스템 국산화 ('21~'23)",
    "• 야전병원 ICT융합 플랫폼 디자인개발 ('21~'25)",
    "",
    "Strengths: 글로벌 Top-tier 수요처와 긴밀한 협력관계",
    "Challenges: 글로벌 다변화 필요, IP 확보 및 R&D 역량 제고 필수",
])

# Key achievements summary
add_card(slide, Inches(8.6), Inches(2.2), Inches(4.0), Inches(4.8), "핵심 역량 요약", [
    "국내 최초 ISO 컨테이너 개발",
    "아시아 최초 A60 Cabin 개발",
    "국내 최초 해외 ESS 공급",
    "30개국 이상 글로벌 공급 경험",
    "특허 13건 (등록 11 + 출원 2)",
    "해외특허 3건 (미국·중국·PCT)",
    "국가R&D 3건 수행",
    "",
    "⚠ 위기대응 경험:",
    "2024년 고객사 파산 (1억달러 손실)",
    "→ R&D 투자 5억→10억 확대",
    "→ 주주증자 550억 확보",
])

add_page_number(slide, 5)


# ===========================
# SLIDE 6: 글로벌 비전 & 성장로드맵 (C. 성장전략 25점)
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_section_header(slide, "02  연구소 비전 & 글로벌 성장 로드맵", "C. 연구소 성장전략 — 글로벌화 전략 (15점) + R&D 성장전략 (10점)")

# Vision bar
add_shape(slide, Inches(0.6), Inches(2.1), Inches(12.1), Inches(0.7), DARK_BLUE)
add_textbox(slide, Inches(0.8), Inches(2.15), Inches(11.5), Inches(0.6),
    'VISION: "글로벌에너지와 전력인프라의 표준이 되는 Modular Containerized Solution"',
    font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 5-stage roadmap
stages = [
    ("~2026", "ESS Value Chain\n내 기술다각화", "• Data Center\n• Grid-forming 인버터\n• SKID 솔루션", ACCENT_BLUE),
    ("2027", "전력기기용\n컨테이너 솔루션", "• 전력기기 통합\n• Grid-forming 양산\n• SKID 상용화", RGBColor(0x00, 0x89, 0x7B)),
    ("2028", "H/W Total Solution\n(Plug & Play)", "• SKID Solution 통합\n• 전력변환/냉각 통합\n• Plug & Play", RGBColor(0x5C, 0x6B, 0xC0)),
    ("2029", "응용분야 확대를\n통한 사업다각화", "• 데이터센터·선박\n• 마이크로그리드\n• 수소연료전지", ORANGE),
    ("2030+", "서비스화", "• O&M 서비스\n• 에너지 트레이딩\n• 플랫폼 비즈니스", RGBColor(0xE9, 0x1E, 0x63)),
]

for i, (year, title, desc, color) in enumerate(stages):
    x = Inches(0.6) + Inches(i * 2.5)
    # Year badge
    add_shape(slide, x, Inches(3.2), Inches(2.3), Inches(0.5), color)
    add_textbox(slide, x, Inches(3.22), Inches(2.3), Inches(0.45), year, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, x + Inches(0.1), Inches(3.85), Inches(2.1), Inches(0.8), title, font_size=12, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    # Description
    add_textbox(slide, x + Inches(0.1), Inches(4.7), Inches(2.1), Inches(1.5), desc, font_size=10, color=MID_GRAY, alignment=PP_ALIGN.LEFT)

# Arrow connectors
for i in range(4):
    x = Inches(2.9) + Inches(i * 2.5)
    add_textbox(slide, x, Inches(3.3), Inches(0.3), Inches(0.4), "→", font_size=20, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Bottom note
add_shape(slide, Inches(0.6), Inches(6.3), Inches(12.1), Inches(0.6), LIGHT_BLUE)
add_textbox(slide, Inches(0.8), Inches(6.35), Inches(11.5), Inches(0.5),
    "GATC 과제 = 성장로드맵 2~3단계 가속화 + 4~5단계 도약의 핵심 전환기술 확보",
    font_size=13, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 6)


# ===========================
# SLIDE 7: 글로벌 목표시장 & 진출전략 (C. 글로벌화)
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_section_header(slide, "02  글로벌 목표시장 & 진출 전략", "글로벌 다변화 — 특정 수요처 의존도 탈피")

# 3 market cards
markets = [
    ("미국 시장", ACCENT_BLUE, [
        "• Fluence Energy — 글로벌 ESS 1위",
        "• PFE/FEOC 공급망 리스크 없는 점 강조",
        "• IRA Domestic Contents 인증 대응",
        "• 현지 법인 운영으로 대응 역량 확보",
    ]),
    ("유럽 시장", GREEN, [
        "• Hitachi 유럽 — 유럽시장 선도",
        "• 친환경·Cybersecurity 강조",
        "• CE 인증 대응",
        "• LOI: Hitachi, PTC Solar, NexGen",
    ]),
    ("아시아·기타", ORANGE, [
        "• 호주, 남미, 동남아",
        "• 재생에너지 확대 빠른 시장",
        "• 현지 파트너십 구축",
        "• 베트남 법인 운영 중",
    ]),
]

for i, (title, color, items) in enumerate(markets):
    x = Inches(0.6) + Inches(i * 4.1)
    add_card(slide, x, Inches(2.2), Inches(3.8), Inches(3.0), title, items, title_color=color)

# Bottom: Localization strategy
add_card(slide, Inches(0.6), Inches(5.5), Inches(12.1), Inches(1.5), "현지화 전략", [
    "• 시장별 규제 DB 구축  • 인증 전담팀 구성  • 설계단계 국제표준 반영  • 현지 규제전문가 네트워크 구축",
    "• 주 1~2회 정기 화상회의  • 1~2개월 1회 베트남·한국 대면미팅  • 분기/반기별 미국 본사 대면미팅",
], title_color=DARK_BLUE)

add_page_number(slide, 7)


# ===========================
# SLIDE 8: 핵심활동 & 세부전략 (C. 성장전략)
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_section_header(slide, "02  핵심활동 세부전략", "글로벌 연구협력·인력파견·세미나·마케팅 연계")

# Strategy cards - 2x2
strategies = [
    ("글로벌 연구협력", [
        "• 탐색→공동연구→전략협력 3단계 운영",
        "• 공동과제 3건/년 목표",
        "• 포트폴리오+표준계약+성과관리",
    ]),
    ("기술인력 해외파견", [
        "• 복지/견학이 아닌 성과중심 기술이전",
        "• 공동목표-공동거버넌스-공동KPI",
        "• 핵심인력 5명/년 파견 목표",
    ]),
    ("해외전문가 초청세미나", [
        "• 행사가 아닌 리드 생성 엔진",
        "• Scoping Call→과제후보→NDA→공동연구 SOW",
        "• 연 4~12회, 리드전환율 20~40%",
    ]),
    ("글로벌 마케팅 연계", [
        "• 연구소 성과를 마케팅 자료로 즉각 활용",
        "• 기술백서·Webinar·Podcast 발간",
        "• RE+·Smarter E Europe 합동 참가",
    ]),
]

for i, (title, items) in enumerate(strategies):
    col = i % 2
    row = i // 2
    x = Inches(0.6) + Inches(col * 6.3)
    y = Inches(2.2) + Inches(row * 2.6)
    add_card(slide, x, y, Inches(6.0), Inches(2.3), title, items, title_color=ACCENT_BLUE if col == 0 else GREEN)

add_page_number(slide, 8)


# ===========================
# SLIDE 9: R&D 투자·인력·인프라 (D. 투자타당성 10점)
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_section_header(slide, "03  R&D 투자·인력·인프라 계획", "D. 투자타당성 및 고용창출 — R&D 인력 (5점) + 글로벌 인력 (5점)")

# KPI Table (as cards)
kpi_data = [
    ("지표", "2026", "2027", "2028", "2029", "2030"),
    ("R&D 매출증가율", "-", "100%", "275%", "409%", "334%"),
    ("신규판로개척(건)", "1", "3", "5", "10", "15"),
    ("글로벌협력실전(건)", "1", "2", "3", "5", "7"),
    ("R&D 연구인력(명)", "12", "15", "18", "21", "25"),
    ("중앙연구소인원(명)", "25", "28", "31", "34", "38"),
    ("특허출원(건)", "2", "3", "3", "4", "4"),
]

# Table header
col_widths = [Inches(2.8), Inches(1.6), Inches(1.6), Inches(1.6), Inches(1.6), Inches(1.6)]
x_start = Inches(0.6)
y_start = Inches(2.2)

for row_idx, row in enumerate(kpi_data):
    x = x_start
    y = y_start + Inches(row_idx * 0.52)
    for col_idx, cell in enumerate(row):
        w = col_widths[col_idx]
        if row_idx == 0:
            bg = DARK_BLUE
            txt_color = WHITE
            is_bold = True
        else:
            bg = WHITE if row_idx % 2 == 1 else LIGHT_BLUE
            txt_color = DARK_GRAY
            is_bold = col_idx == 0
        add_shape(slide, x, y, w, Inches(0.48), bg)
        add_textbox(slide, x, y + Inches(0.05), w, Inches(0.38), cell,
            font_size=12, color=txt_color, bold=is_bold, alignment=PP_ALIGN.CENTER)
        x += w

# Bottom: Growth rate note
add_textbox(slide, Inches(0.6), Inches(6.0), Inches(12), Inches(0.4),
    "연평균 성장률: R&D 매출 49.5%  |  신규판로 71.0%  |  글로벌협력 62.7%  |  연구인력 20.1%",
    font_size=13, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# Human resource plan
add_card(slide, Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.8), "글로벌 인력확보 전략", [
    "• 글로벌 유수 전력기기업체 연구소 출신 인재 영입  • 박사급 전문인력 선제채용  • 부산산학협력센터 인턴→채용 연계",
], title_color=GREEN)

add_page_number(slide, 9)


# ===========================
# SLIDE 10: 협력 네트워크
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_section_header(slide, "04  국내외 협력 네트워크", "R&D 자문·시험평가·공동연구·실증 협력")

# Domestic network
add_card(slide, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.8), "국내 네트워크", [
    "R&D 자문·시험평가:",
    "  한국전기연구원 — ESS·GFM 시험평가",
    "  한국생산기술연구원 — 제조기술 자문",
    "  한국스마트그리드협회 — 신재생전력망 자문",
    "",
    "산학협력:",
    "  서울대학교 — GFM 핵심기술 중장기연구",
    "  성균관대학교 — 전력시스템 연구협력",
    "  동아대학교 — 인재육성·산학협력",
])

# Overseas network
add_card(slide, Inches(6.8), Inches(2.2), Inches(5.8), Inches(4.8), "해외 네트워크 & 수요처", [
    "공동연구·실증:",
    "  EPC Power — GFM 인버터 공동개발",
    "  Fluence — 글로벌 실증·사업화",
    "  CEC — 시험평가·인증 협업",
    "",
    "수요처 네트워크:",
    "  System Integrator: Fluence, Energy Vault, Qcells",
    "  전력기기업체: ABB, Hitachi, On Energy",
    "  시공/EPC: McCarthy, DPR, NexGen",
    "",
    "LOI 체결: Hitachi, Fluence, PTC Solar, NexGen",
])

add_page_number(slide, 10)


# ===========================
# SLIDE 11: 부합성 & 경영진 의지 (C. 성장전략)
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, LIGHT_GRAY)
add_section_header(slide, "05  기업-연구소 부합성 & 경영진 의지", "연구소가 지원조직이 아닌 성장 플랫폼")

# Alignment cards
add_card(slide, Inches(0.6), Inches(2.2), Inches(3.8), Inches(2.2), "시장 다변화", [
    "• ESS 컨테이너 편중 → 데이터센터·선박·마이크로그리드",
    "• 특정 수요처 의존도 탈피 → 다변화된 파이프라인",
])

add_card(slide, Inches(4.7), Inches(2.2), Inches(3.8), Inches(2.2), "고객 종속 리스크 분산", [
    "• Fluence 의존도 → 유럽/아시아 신규 고객 확보",
    "• IP 포트폴리오 강화로 기술 자립성 확보",
])

add_card(slide, Inches(8.8), Inches(2.2), Inches(3.8), Inches(2.2), "글로벌 경쟁력 확보", [
    '• "연구소의 성장이 곧 기업의 성장"',
    "• 연구소 성과를 마케팅·영업에 즉각 활용",
])

# Management commitment
add_card(slide, Inches(0.6), Inches(4.7), Inches(12.1), Inches(2.5), "경영진 연구소 육성 의지", [
    "✓ 매출 무관 R&D 연 5억+ 하한선 유지                              ✓ 데이터센터·스키드·GFM인버터 우선투자",
    "✓ 부산산학협력센터 운영투자                                       ✓ Gate 운영: 아이템→PoC→검증→제품화",
    "✓ 표준화·모듈화·문서화 → 플랫폼 자산화                         ✓ 핵심인력 선제채용 (전력변환·열관리·인증)",
    "✓ 인턴→채용 파이프라인, 리텐션 강화",
    "",
    "연구소 3년내 플랫폼 승격  |  IP 생산기지  |  글로벌 게이트키퍼",
], title_color=DARK_BLUE)

add_page_number(slide, 11)


# ===========================
# SLIDE 12: 핵심요약 & 마무리
# ===========================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_BLUE)

add_textbox(slide, Inches(1), Inches(0.5), Inches(11), Inches(0.7),
    "핵심 요약", font_size=36, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

# 3 summary cards
summaries = [
    ("01  연구소 역량", ACCENT_BLUE, [
        "45년 컨테이너 역사",
        "글로벌 1위 ESS 공급",
        "13건 특허 · 인증체계 완비",
        "기술사업화 누적 1조+ 매출",
    ]),
    ("02  성장 계획", GREEN, [
        "5단계 글로벌 성장로드맵",
        "R&D 투자 2배 확대 (5억→10억)",
        "25명 전문인력 확보 목표",
        "연도별 KPI 체계 수립",
    ]),
    ("03  GATC = 핵심 전환점", ORANGE, [
        "기업전략과 완벽 부합",
        "경영진 R&D 의지 확고",
        "성장로드맵 2~3단계 가속화",
        "전력기기 통합 플랫폼 도약",
    ]),
]

for i, (title, color, items) in enumerate(summaries):
    x = Inches(0.6) + Inches(i * 4.2)
    add_card(slide, x, Inches(1.8), Inches(3.9), Inches(3.5), title, items, title_color=color, bg_color=WHITE)

# Thank you
add_shape(slide, Inches(0), Inches(5.8), prs.slide_width, Inches(0.04), ACCENT_BLUE)
add_textbox(slide, Inches(1), Inches(6.0), Inches(11), Inches(0.8),
    "경청해 주셔서 감사합니다", font_size=28, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)

add_page_number(slide, 12)


# ===== SAVE =====
output_path = "/root/.openclaw/workspace/GATC_성장전략계획서_에이스엔지니어링_10분.pptx"
prs.save(output_path)
print(f"✅ PPT saved: {output_path}")
print(f"   Slides: {len(prs.slides)}")
